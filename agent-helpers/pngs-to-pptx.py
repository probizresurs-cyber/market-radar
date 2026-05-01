#!/usr/bin/env python3
"""
pngs-to-pptx.py — Combine a directory of PNG slide screenshots into a .pptx file.

Each PNG becomes a full-bleed background of a 16:9 slide (13.333" x 7.5").
The slides are NOT editable as text — the user gets a beautiful HTML-rendered
deck packaged as PPTX for compatibility with PowerPoint / Keynote / Google Slides.

Usage:
    python3 pngs-to-pptx.py <input-dir> <output.pptx>
"""

import sys
import os
import glob
from pptx import Presentation
from pptx.util import Inches


def main():
    if len(sys.argv) < 3:
        print("Usage: python3 pngs-to-pptx.py <input-dir> <output.pptx>", file=sys.stderr)
        sys.exit(1)

    input_dir = sys.argv[1]
    output_pptx = sys.argv[2]

    pngs = sorted(glob.glob(os.path.join(input_dir, "slide-*.png")))
    if not pngs:
        print(f"No slide-*.png files found in {input_dir}", file=sys.stderr)
        sys.exit(1)

    # 1920x1080 at 144dpi → 13.333" x 7.5" (16:9 widescreen)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]  # blank layout

    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            png,
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(output_pptx)
    print(f"Saved {output_pptx} with {len(pngs)} slides")


if __name__ == "__main__":
    main()
